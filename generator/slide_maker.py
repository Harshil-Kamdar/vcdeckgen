from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

def build_presentation(slides, charts, images):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Intro slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "🚀 Startup Pitch Deck"
    slide.placeholders[1].text = "Auto-generated with VCDeckGen"

    # Content slides
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = slide_data['title']
        content_shape = slide.placeholders[1]
        content_shape.text = slide_data['content']

        # Add chart if exists
        if charts[i]:
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide.shapes.add_picture(charts[i], left, top, height=height)

    # Save deck
    prs.save("output/final_presentation.pptx")
